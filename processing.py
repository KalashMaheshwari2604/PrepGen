try:
    import fitz  # PyMuPDF
except ImportError:
    # Fallback if fitz has issues
    try:
        import pymupdf as fitz
    except ImportError:
        fitz = None

import docx
from pptx import Presentation
import os


def extract_text(file_path: str) -> str:
    """
    Extracts text from various file types with improved logic for PPTX.
    """
    _, ext = os.path.splitext(file_path)
    ext = ext.lower()
    text = ""

    try:
        if ext == '.pdf':
            # Use pymupdf directly (PyMuPDF package)
            import pymupdf
            doc = pymupdf.open(file_path)
            try:
                for page in doc:
                    text += page.get_text() + "\n"
            finally:
                doc.close()
        elif ext == '.docx':
            doc = docx.Document(file_path)
            for para in doc.paragraphs:
                text += para.text + "\n"
        elif ext == '.pptx':
            # --- Smarter PPTX Extraction Logic ---
            prs = Presentation(file_path)
            for slide in prs.slides:
                slide_texts = []
                # Try to get the title first
                if slide.shapes.title:
                    slide_texts.append(slide.shapes.title.text)
                
                # Now get the main content placeholders
                for shape in slide.shapes:
                    if shape.has_text_frame and shape.is_placeholder:
                        # Check for body, content, and subtitle placeholders
                        if shape.placeholder_format.idx in [1, 13, 14, 15, 16]:
                            slide_texts.append(shape.text)
                
                text += "\n".join(slide_texts) + "\n\n"
        elif ext == '.txt':
            with open(file_path, 'r', encoding='utf-8') as f:
                text = f.read()
        else:
            print(f"Unsupported file type: {ext}")
            return ""

    except Exception as e:
        print(f"Error reading {file_path}: {e}")
        return ""

    cleaned_text = "\n".join(line.strip() for line in text.splitlines() if line.strip())
    return cleaned_text

def chunk_text(text: str, chunk_size: int = 768, chunk_overlap: int = 100) -> list[str]:
    """
    Splits a long text into smaller, overlapping chunks.

    Args:
        text (str): The input text.
        chunk_size (int): The maximum size of each chunk (in characters).
        chunk_overlap (int): The number of characters to overlap between chunks.

    Returns:
        list[str]: A list of text chunks.
    """
    if not text:
        return []

    chunks = []
    start = 0
    while start < len(text):
        end = start + chunk_size
        chunks.append(text[start:end])
        # Move the start pointer back by the overlap amount
        start += chunk_size - chunk_overlap

    return chunks

# --- This is our test block ---
if __name__ == '__main__':
    # 1. Define the path to your sample file
    # Make sure 'sample.txt' is in the same folder as this script
    test_file_path = 'sample.txt'

    # 2. Extract text from the file
    print(f"--- Extracting text from {test_file_path} ---")
    extracted_text = extract_text(test_file_path)
    if extracted_text:
        print(f"Successfully extracted {len(extracted_text)} characters.")
        print("First 100 characters:", extracted_text[:100], "...")
        print("\n" + "="*50 + "\n")

        # 3. Chunk the extracted text
        print("--- Chunking the text ---")
        text_chunks = chunk_text(extracted_text)
        if text_chunks:
            print(f"Split text into {len(text_chunks)} chunks.")
            print("--- First Chunk ---")
            print(text_chunks[0])
            print("--- Second Chunk (showing overlap) ---")
            print(text_chunks[1] if len(text_chunks) > 1 else "No second chunk.")
    else:
        print("Text extraction failed.")